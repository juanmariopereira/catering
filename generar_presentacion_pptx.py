#!/usr/bin/env python
"""
Genera la presentación PowerPoint del sistema de catering (resumen ejecutivo).
Ejecutar: python generar_presentacion_pptx.py
Salida: presentacion_catering.pptx en la raíz del proyecto.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor

# Layouts: 0 = title, 1 = title + body
TITLE_SLIDE = 0
TITLE_AND_BODY = 1


def add_bullet_slide(prs, title, bullets, layout_idx=1):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(14)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1 - Portada
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    slide.shapes.title.text = "Sistema de Gestión de Catering"
    try:
        slide.placeholders[1].text = "Resumen ejecutivo para interesados"
    except IndexError:
        pass

    # Slide 2 - Qué resuelve
    add_bullet_slide(prs, "¿Qué resuelve el sistema?", [
        "Un solo sistema para todo el ciclo del servicio",
        "Centraliza clientes, contratos, planes y condiciones de entrega",
        "Planifica menús por fecha y plan con sustituciones por cliente",
        "Organiza producción en cocina y rutas de entrega",
        "Genera previsiones de compra a partir de la planificación",
        "Gestiona cobros, pagos y vencimientos con dashboard de cobranza",
    ])

    # Slide 3 - Fortalezas 1
    add_bullet_slide(prs, "Fortalezas principales (1/2)", [
        "Contratos inteligentes: pausas (vacaciones) y feriados recalculan automáticamente la vigencia del contrato",
        "Días extra con trazabilidad: extensión de vigencia del contrato y del último cobro con motivo y registro de auditoría",
        "Planificación por fecha y plan: menús por momento del día (desayuno, comida, cena, etc.) con sustituciones por cliente",
        "Previsión de compra automática: cantidades de ingredientes calculadas desde los menús planificados en un rango de fechas",
    ])

    # Slide 4 - Fortalezas 2
    add_bullet_slide(prs, "Fortalezas principales (2/2)", [
        "Cocina por fecha: resumen de platos, cantidades, cambios por cliente y notas para el equipo",
        "Rutas y entregas: asignación de clientes a entregadores, códigos de entrega y hoja de ruta imprimible",
        "Cobranza integrada: cobros, pagos (por defecto QR), estados (pendiente/vencida/pagada) y dashboard de cobranza",
        "Dashboard principal: indicadores en tiempo real, gráficos de pagos y previsión de cobros, acciones rápidas",
    ])

    # Slide 5 - Tres pilares
    add_bullet_slide(prs, "Tres pilares del sistema", [
        "Catálogos: Clientes, Planes, Recetas, Ingredientes, Dietas, Feriados, Entregadores",
        "Producción: Planificación de menús, Cocina por fecha, Composición por fecha, Entregas y rutas, Previsiones de compra",
        "Finanzas: Contratos, Cobros, Pagos, Dashboard de cobranza",
    ])

    # Slide 6 - Feriados
    add_bullet_slide(prs, "Feriados integrados en la operación", [
        "Un solo calendario de feriados para todo el sistema",
        "Contratos: las pausas que incluyen un feriado no cuentan ese día; la fecha de fin se recalcula solo",
        "Entregas: en feriados no se ofrecen contratos para asignar a ruta",
        "Planificación y cocina: aviso de feriado en las vistas por fecha",
    ])

    # Slide 7 - Cierre
    add_bullet_slide(prs, "Conclusión", [
        "Sistema integral: desde el contrato hasta la entrega y el cobro",
        "Automatización donde importa: vigencia, feriados, previsiones de compra, indicadores",
        "Trazabilidad: días extra, pagos y estados de cobro",
        "Listo para operar día a día y para presentar a interesados",
    ])

    # Slide 8 - Próximos pasos
    add_bullet_slide(prs, "Próximos pasos", [
        "Demo en vivo del sistema",
        "Capacitación a usuarios",
        "Definición de acceso piloto (si aplica)",
    ])

    out_path = "presentacion_catering.pptx"
    prs.save(out_path)
    print(f"Presentación guardada: {out_path}")


if __name__ == "__main__":
    main()
